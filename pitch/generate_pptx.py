#!/usr/bin/env python3
"""Generate Auteursportaal CEO pitch — Noordhoff template design system."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

TPL = '/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/Template richtlijnen - 2026.pptx'
SHOTS = '/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/screenshots'

# ── Load template ──
prs = Presentation(TPL)

# Remove all template slides (keep layouts only)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

W = prs.slide_width
H = prs.slide_height

# ══════════════════════════════════════════════
# NOORDHOFF DESIGN SYSTEM — Colors & Typography
# ══════════════════════════════════════════════

TEAL = RGBColor(0, 122, 94)         # #007A5E — primary brand
BLACK = RGBColor(0, 0, 0)           # titles
BODY = RGBColor(41, 41, 41)         # #292929 — body text
SUBTITLE_CLR = RGBColor(50, 50, 50) # #323232 — subtitles
ORANGE = RGBColor(255, 112, 67)     # #FF7043 — accent/negative
WHITE = RGBColor(255, 255, 255)
WHITE_SHAPE = RGBColor(254, 255, 255)  # #FEFFFF from template
LIGHT_GREY = RGBColor(180, 180, 180)
BORDER_CLR = RGBColor(220, 220, 220)

# Layout references
LY_TITLE = prs.slide_layouts[0]       # Title Slide
LY_TITLE_ONLY = prs.slide_layouts[4]  # Title Only
LY_GREEN = prs.slide_layouts[5]       # Title Only Green
LY_BLANK = prs.slide_layouts[8]       # Blank
LY_BLANK_GREEN = prs.slide_layouts[9] # Blank Green

FONT = 'Roboto Light'
FONT_HAND = 'Caveat'


# ══════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════

def add_text(slide, l, t, w, h, text, sz=14, color=BODY, align=PP_ALIGN.LEFT,
             font=FONT, spacing=None, bold=False, line_spacing=None):
    """Add a text box. Supports multi-line via \\n."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = font
        p.font.bold = bold
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return box


def add_title(slide, text, t=Inches(0.5)):
    """Slide title — black, ~40pt, left-aligned at x=0.9"."""
    add_text(slide, Inches(0.9), t, Inches(11), Inches(0.8),
             text, sz=40, color=BLACK, font=FONT)


def add_subtitle_label(slide, text, t=Inches(0.3)):
    """Small category label above title — teal, uppercase, 11pt."""
    add_text(slide, Inches(0.9), t, Inches(6), Inches(0.3),
             text.upper(), sz=11, color=TEAL, font=FONT)


def add_teal_square(slide, l, t, size, opacity_adjust=0):
    """Decorative teal rounded square — visual anchor."""
    sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t,
                                 Inches(size), Inches(size))
    sq.fill.solid()
    sq.fill.fore_color.rgb = TEAL
    sq.line.fill.background()
    # Gentle rounding
    sq.adjustments[0] = 0.08
    return sq


def add_marker_dot(slide, l, t, color=TEAL, size=0.12):
    """Small teal rounded square used as bullet marker."""
    m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t,
                                Inches(size), Inches(size))
    m.fill.solid()
    m.fill.fore_color.rgb = color
    m.line.fill.background()
    m.adjustments[0] = 0.25
    return m


def add_thin_line(slide, l, t, w, color=BORDER_CLR):
    """Thin horizontal separator line."""
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def add_screenshot(slide, filename, l, t, w, label=None):
    """Add a screenshot image with optional label above it."""
    path = os.path.join(SHOTS, filename)
    if label:
        add_text(slide, l, t - Inches(0.25), Inches(3), Inches(0.25),
                 label, sz=9, color=TEAL, font=FONT)
    if os.path.exists(path):
        # Add subtle border/shadow effect via a background rectangle
        border_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, l - Inches(0.02), t - Inches(0.02),
            w + Inches(0.04), Inches(0.04)  # placeholder height, will be covered
        )
        border_rect.fill.solid()
        border_rect.fill.fore_color.rgb = BORDER_CLR
        border_rect.line.fill.background()
        # Remove the border rect — we don't actually need it as picture has own bounds
        # Just add the picture directly
        slide.shapes.add_picture(path, l, t, w)


# ══════════════════════════════════════════════
# SLIDE 1 — TITLE (split: screenshot left, green right)
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

# Left half: screenshot as teaser (faded)
path = os.path.join(SHOTS, '02_start.png')
if os.path.exists(path):
    s.shapes.add_picture(path, Inches(-0.5), Inches(-0.2), Inches(8.0))

# Right half: teal block
teal_block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(6.8), Inches(0), Inches(6.6), H)
teal_block.fill.solid(); teal_block.fill.fore_color.rgb = TEAL
teal_block.line.fill.background()

# Decorative white pills on teal
for (pl, pt, pw, ph) in [
    (Inches(7.2), Inches(0.5), Inches(2.0), Inches(0.25)),
    (Inches(11.0), Inches(6.6), Inches(2.0), Inches(0.3)),
]:
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pl, pt, pw, ph)
    pill.fill.solid(); pill.fill.fore_color.rgb = WHITE_SHAPE
    pill.line.fill.background(); pill.adjustments[0] = 0.5

# Title text on teal
add_text(s, Inches(7.3), Inches(2.0), Inches(5.5), Inches(1.2),
         'Auteurs-\nportaal', sz=52, color=WHITE, font=FONT)

# Subtitle
add_text(s, Inches(7.3), Inches(3.8), Inches(5.0), Inches(0.8),
         'Het digitale platform voor\nroyalty-afrekeningen, contracten\nen prognoses',
         sz=15, color=RGBColor(200, 235, 225), spacing=6)

# Divider
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(7.3), Inches(5.1), Inches(2.0), Pt(1.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0, 150, 120)
div.line.fill.background()

add_text(s, Inches(7.3), Inches(5.3), Inches(5), Inches(0.4),
         'Noordhoff  ·  Liber  ·  Plantyn', sz=13, color=RGBColor(160, 215, 200))

add_text(s, Inches(7.3), Inches(6.2), Inches(5), Inches(0.3),
         'Patrick Jeeninga  —  Maart 2026', sz=10, color=RGBColor(120, 185, 165))


# ══════════════════════════════════════════════
# SLIDE 2 — CONTEXT: WAAROM NU?
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

# Decorative teal squares
add_teal_square(s, Inches(-0.5), Inches(-0.5), 1.0)
add_teal_square(s, Inches(12.2), Inches(6.0), 1.8)

add_subtitle_label(s, 'Context')
add_title(s, 'Waarom nu?', t=Inches(0.55))

items = [
    ('Aandeelhouders',
     'De aandeelhouders pushen actief op meer digitalisatie.\nDit portaal is een concrete, zichtbare stap die direct waarde levert.'),
    ('Value Creation Plan',
     'Past direct binnen het Value Creation Plan als digitaal\nverbeterproject. Schaalbaar naar Liber en Plantyn.'),
    ('250 auteurs onzeker',
     '10% van alle auteurs belt of mailt jaarlijks met vragen\nover hun afrekening. Gebrek aan transparantie.'),
    ('Concurrentiepositie',
     'Auteurs vergelijken hun ervaring met andere uitgevers.\nEen modern portaal is een retentie-instrument.'),
]

for i, (title, desc) in enumerate(items):
    col, row = i % 2, i // 2
    l = Inches(0.9 + col * 6.0)
    t = Inches(1.7 + row * 2.6)

    # Marker dot — consistent teal
    add_marker_dot(s, l, t + Inches(0.07), color=TEAL)

    # Title
    add_text(s, l + Inches(0.3), t, Inches(5.2), Inches(0.35),
             title, sz=18, color=BLACK, font=FONT)

    # Description
    add_text(s, l + Inches(0.3), t + Inches(0.5), Inches(5.2), Inches(1.5),
             desc, sz=13, color=BODY, spacing=5)

    # Separator line under each item
    if row == 0:
        add_thin_line(s, l + Inches(0.3), t + Inches(2.1), Inches(5.0))


# ══════════════════════════════════════════════
# SLIDE 3 — PROBLEM: AUTEUR EXPERIENCE
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.0), Inches(-0.4), 1.2)

add_subtitle_label(s, 'Het probleem')
add_title(s, 'Hoe ervaart de auteur het nu?', t=Inches(0.55))

# Outlook-style email mock
email_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2.6))
email_bg.fill.solid(); email_bg.fill.fore_color.rgb = WHITE
email_bg.line.color.rgb = BORDER_CLR; email_bg.line.width = Pt(1); email_bg.adjustments[0] = 0.02

# Email toolbar bar
toolbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.35))
toolbar.fill.solid(); toolbar.fill.fore_color.rgb = RGBColor(243, 243, 243); toolbar.line.fill.background()

add_text(s, Inches(1.8), Inches(1.53), Inches(3), Inches(0.3),
         'Beantwoorden    Doorsturen    Verwijderen', sz=8, color=LIGHT_GREY)

# Email headers
header_y = Inches(2.0)
for label, value in [
    ('Van:', 'M. Hendriks <m.hendriks@gmail.com>'),
    ('Aan:', 'rights@noordhoff.nl'),
    ('Onderwerp:', 'Vraag over afrekening Getal & Ruimte 2024'),
]:
    add_text(s, Inches(1.8), header_y, Inches(1.2), Inches(0.22),
             label, sz=9, color=LIGHT_GREY)
    add_text(s, Inches(2.9), header_y, Inches(8), Inches(0.22),
             value, sz=9, color=BODY)
    header_y += Inches(0.22)

# Email body
add_text(s, Inches(1.8), Inches(2.8), Inches(9.5), Inches(1),
         'Beste Noordhoff,\n\n'
         'Ik heb in maart een brief ontvangen over Getal & Ruimte maar ik begrijp\n'
         'de berekening niet. Kan iemand mij bellen? Ik wil ook graag weten of\n'
         'mijn contract nog loopt en wat de prognose is voor volgend jaar.\n\n'
         'Met vriendelijke groet,\nM. Hendriks',
         sz=10, color=BODY, spacing=2)

# Annotation
add_text(s, Inches(8.5), Inches(1.2), Inches(4), Inches(0.4),
         '1 van 250 per jaar', sz=20, color=BODY, font=FONT_HAND)

add_thin_line(s, Inches(0.9), Inches(4.4), Inches(11.5))

# Three pain points
pains = [
    ('1 brief per jaar',
     'De enige communicatie over royalties is een\njaarlijkse brief in maart. De rest van het jaar: stilte.'),
    ('Geen inzicht',
     'Auteurs kunnen niet zelf opzoeken hoe hun\nmethode presteert. Volledig afhankelijk van Noordhoff.'),
    ('Onzekerheid',
     '250 auteurs nemen jaarlijks contact op\nmet vragen. Dat is 10% van alle auteurs.'),
]

for i, (title, desc) in enumerate(pains):
    l = Inches(0.9 + i * 4.0)

    add_marker_dot(s, l, Inches(4.97), color=TEAL)
    add_text(s, l + Inches(0.3), Inches(4.9), Inches(3.3), Inches(0.35),
             title, sz=16, color=BLACK, font=FONT)
    add_text(s, l + Inches(0.3), Inches(5.4), Inches(3.3), Inches(1.5),
             desc, sz=12, color=BODY, spacing=5)


# ══════════════════════════════════════════════
# SLIDE 4 — PROBLEM: COSTS
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(5.8), 1.5)

add_subtitle_label(s, 'Het probleem')
add_title(s, 'Wat kost het Noordhoff nu?', t=Inches(0.55))

# Four stat blocks — clean, no card backgrounds
stats = [
    ('\u20AC10.000', 'Postkosten per jaar\n(3 organisaties)'),
    ('32 uur', 'Brieven inpakken\n(2 man \u00D7 2 dagen)'),
    ('250', 'Auteurs met vragen\nper jaar'),
    ('0', 'Digitaal inzicht\nvoor auteurs'),
]
for i, (value, label) in enumerate(stats):
    l = Inches(0.9 + i * 3.05)
    # Large value
    add_text(s, l, Inches(1.7), Inches(2.7), Inches(0.7),
             value, sz=36, color=TEAL, font=FONT, align=PP_ALIGN.LEFT)
    # Label below
    add_text(s, l, Inches(2.4), Inches(2.7), Inches(0.6),
             label, sz=11, color=LIGHT_GREY, spacing=3)

add_thin_line(s, Inches(0.9), Inches(3.3), Inches(11.5))

# Hidden costs list
add_text(s, Inches(0.9), Inches(3.6), Inches(11), Inches(0.35),
         'Verborgen kosten die niet op de factuur staan:', sz=14, color=BLACK, font=FONT)

hidden = [
    ('Afhandeling auteursvragen', '250 vragen \u00D7 gem. 30 min = 125 uur per jaar', '~\u20AC5.000'),
    ('Fysieke verwerking', '2 medewerkers \u00D7 2 dagen printen, vouwen, inpakken', '~\u20AC2.500'),
    ('Reputatierisico', 'Auteurs ervaren Noordhoff als ouderwets en niet-transparant', '\u2014'),
    ('Auteursverlies', 'Frustratie leidt tot overstap naar concurrent', '\u2014'),
]
for i, (item, detail, cost) in enumerate(hidden):
    top = Inches(4.2 + i * 0.55)

    add_marker_dot(s, Inches(1.0), top + Inches(0.08), color=ORANGE, size=0.1)
    add_text(s, Inches(1.3), top, Inches(3.2), Inches(0.4),
             item, sz=12, color=BLACK, font=FONT)
    add_text(s, Inches(4.6), top, Inches(5.2), Inches(0.4),
             detail, sz=11, color=LIGHT_GREY)
    add_text(s, Inches(10.5), top, Inches(1.8), Inches(0.4),
             cost, sz=12, color=ORANGE, align=PP_ALIGN.RIGHT, font=FONT)

# Total line
add_thin_line(s, Inches(0.9), Inches(6.5), Inches(11.5), color=TEAL)
add_text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.3),
         'Totale kosten huidige situatie', sz=13, color=TEAL, font=FONT)
add_text(s, Inches(10.5), Inches(6.6), Inches(1.8), Inches(0.3),
         '~\u20AC18.000/jr', sz=13, color=TEAL, align=PP_ALIGN.RIGHT, font=FONT)


# ══════════════════════════════════════════════
# SLIDE 5 — TRANSFORMATION: OLD VS NEW
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(-0.5), 1.0)

add_subtitle_label(s, 'De transformatie')
add_title(s, 'Van papier naar portaal', t=Inches(0.55))

# LEFT — OUD label block
old_label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.9), Inches(1.6), Inches(1.0), Inches(0.35))
old_label.fill.solid(); old_label.fill.fore_color.rgb = ORANGE
old_label.line.fill.background(); old_label.adjustments[0] = 0.4
add_text(s, Inches(0.9), Inches(1.62), Inches(1.0), Inches(0.3),
         'OUD', sz=10, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(2.1), Inches(1.6), Inches(4), Inches(0.35),
         'Fysieke post', sz=18, color=BODY, font=FONT)

old_items = [
    '1× per jaar een royalty-brief',
    'Geen tussentijds inzicht',
    'Vragen via email en telefoon',
    'Contracten in archiefkast',
    'Geen prognose mogelijk',
    'Gegevenswijziging per formulier',
    'Brieven handmatig inpakken',
    'Geen inzicht in auteuractiviteit',
]
for j, item in enumerate(old_items):
    add_text(s, Inches(1.2), Inches(2.3 + j * 0.45), Inches(4.8), Inches(0.35),
             f'×   {item}', sz=12, color=LIGHT_GREY)

# Central arrow
arrow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.1), Inches(3.5), Inches(1.1), Inches(1.1))
arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL; arrow.line.fill.background()
add_text(s, Inches(6.1), Inches(3.6), Inches(1.1), Inches(0.9),
         '→', sz=32, color=WHITE, align=PP_ALIGN.CENTER)

# RIGHT — NIEUW label block
new_label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(7.5), Inches(1.6), Inches(1.2), Inches(0.35))
new_label.fill.solid(); new_label.fill.fore_color.rgb = TEAL
new_label.line.fill.background(); new_label.adjustments[0] = 0.4
add_text(s, Inches(7.5), Inches(1.62), Inches(1.2), Inches(0.3),
         'NIEUW', sz=10, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, Inches(8.9), Inches(1.6), Inches(4), Inches(0.35),
         'Het Auteursportaal', sz=18, color=TEAL, font=FONT)

new_items = [
    '24/7 inzicht in afrekeningen',
    'Realtime royalty-overzicht per jaar',
    'FAQ en rondleiding in het portaal',
    'Contracten direct inzien en downloaden',
    'Prognose met min/max range',
    'Wijzigingen digitaal aanvragen',
    'Bulk upload + automatisch versturen',
    'Activiteiten feed voor admin',
]
for j, item in enumerate(new_items):
    add_marker_dot(s, Inches(7.8), Inches(2.37 + j * 0.45), color=TEAL, size=0.08)
    add_text(s, Inches(8.05), Inches(2.3 + j * 0.45), Inches(4.3), Inches(0.35),
             item, sz=12, color=BODY)

# Vertical separator line
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(6.65), Inches(1.5), Pt(1), Inches(5.0))
sep.fill.solid()
sep.fill.fore_color.rgb = BORDER_CLR
sep.line.fill.background()


# ══════════════════════════════════════════════
# SLIDE 6 — SOLUTION: TWO DASHBOARDS
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(-0.5), Inches(-0.5), 1.0)
add_teal_square(s, Inches(12.5), Inches(6.2), 1.2)

add_subtitle_label(s, 'De oplossing')
add_title(s, 'E\u00E9n portaal, twee dashboards', t=Inches(0.55))

# LEFT — Auteur dashboard
add_marker_dot(s, Inches(0.9), Inches(1.77), color=TEAL, size=0.14)
add_text(s, Inches(1.2), Inches(1.7), Inches(4.5), Inches(0.35),
         'Auteursdashboard', sz=18, color=TEAL, font=FONT)

auteur_features = [
    'Jaaroverzicht met royalty-totalen',
    'Afrekeningen zoeken, filteren, downloaden',
    'Contracten inzien met PDF preview',
    'Prognose komend jaar (min\u2014max)',
    'Declaraties indienen met PDF upload',
    'Persoonlijke gegevens beheren',
    'Interactieve royalty-grafiek per type',
    'FAQ, guided tour, command palette (\u2318K)',
]
for j, f in enumerate(auteur_features):
    add_text(s, Inches(1.2), Inches(2.3 + j * 0.45), Inches(5), Inches(0.35),
             f'\u00B7   {f}', sz=12, color=BODY)

# Vertical separator
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(6.5), Inches(1.7), Pt(1), Inches(5.0))
sep.fill.solid()
sep.fill.fore_color.rgb = BORDER_CLR
sep.line.fill.background()

# RIGHT — Admin dashboard
add_marker_dot(s, Inches(7.0), Inches(1.77), color=TEAL, size=0.14)
add_text(s, Inches(7.3), Inches(1.7), Inches(4.5), Inches(0.35),
         'Admin Dashboard', sz=18, color=BLACK, font=FONT)

admin_features = [
    'Overzicht van alle auteurs met zoekfunctie',
    'Per auteur: gegevens, contracten, afrekeningen',
    'Wijzigingsverzoeken goedkeuren/afwijzen',
    'Bulk PDF upload voor afrekeningen',
    'CSV import voor nieuwe auteurs',
    'Evenementen en nieuws beheren',
    'Activiteiten feed (logins, wijzigingen)',
    'E-mail notificatie instellingen',
]
for j, f in enumerate(admin_features):
    add_text(s, Inches(7.3), Inches(2.3 + j * 0.45), Inches(5), Inches(0.35),
             f'\u00B7   {f}', sz=12, color=BODY)


# ══════════════════════════════════════════════
# SLIDE 7 — SCREENSHOT: AUTEUR DASHBOARD
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(-0.4), 1.0)

add_subtitle_label(s, 'Het portaal')
add_title(s, 'Wat de auteur ziet', t=Inches(0.55))

# Screenshot left (large)
path = os.path.join(SHOTS, '02_start.png')
if os.path.exists(path):
    # Subtle shadow rectangle behind screenshot
    shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.85), Inches(1.65), Inches(7.6), Inches(5.2))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(235, 235, 235)
    shadow.line.fill.background(); shadow.adjustments[0] = 0.015
    s.shapes.add_picture(path, Inches(0.9), Inches(1.7), Inches(7.5))

# Explanation text right
add_text(s, Inches(8.8), Inches(1.7), Inches(3.8), Inches(0.35),
         'Persoonlijk dashboard', sz=16, color=BLACK, font=FONT)

add_thin_line(s, Inches(8.8), Inches(2.2), Inches(3.5), color=TEAL)

features = [
    ('Jaaroverzicht', 'Totaal uitgekeerd, laatste betaling\nen verwachte royalties in één kaart.'),
    ('Royalty-grafiek', 'Interactieve grafiek per jaar met\nuitsplitsing per type (royalties,\nnevenrechten, foreign rights).'),
    ('Afrekeningen', 'Zoeken, filteren, PDF preview\nen download. CSV-export.'),
    ('Contracten', 'Alle contracten inzien met\nPDF preview en download.'),
    ('Prognose', 'Verwachte royalties komend jaar\nmet min/max range.'),
]
for i, (title, desc) in enumerate(features):
    y = Inches(2.5 + i * 0.9)
    add_marker_dot(s, Inches(8.8), y + Inches(0.05), color=TEAL, size=0.08)
    add_text(s, Inches(9.05), y, Inches(3.5), Inches(0.25),
             title, sz=11, color=BLACK)
    add_text(s, Inches(9.05), y + Inches(0.28), Inches(3.5), Inches(0.6),
             desc, sz=9, color=LIGHT_GREY, spacing=2)


# ══════════════════════════════════════════════
# SLIDE 8 — SCREENSHOT: ADMIN DASHBOARD
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(-0.5), Inches(-0.4), 1.0)

add_subtitle_label(s, 'Het portaal')
add_title(s, 'Wat de beheerder ziet', t=Inches(0.55))

# Screenshot right (large)
path = os.path.join(SHOTS, '10_admin.png')
if os.path.exists(path):
    shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.85), Inches(1.65), Inches(7.6), Inches(5.2))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(235, 235, 235)
    shadow.line.fill.background(); shadow.adjustments[0] = 0.015
    s.shapes.add_picture(path, Inches(4.9), Inches(1.7), Inches(7.5))

# Explanation text left
add_text(s, Inches(0.9), Inches(1.7), Inches(3.6), Inches(0.35),
         'Admin dashboard', sz=16, color=BLACK, font=FONT)

add_thin_line(s, Inches(0.9), Inches(2.2), Inches(3.3), color=TEAL)

admin_feats = [
    ('Auteursbeheer', '9 auteurs met zoekfunctie.\nPer auteur: gegevens, contracten,\nafrekeningen, prognose.'),
    ('Wijzigingen', 'Verzoeken van auteurs goedkeuren\nof afwijzen. Volledige audit trail.'),
    ('Bulk import', 'CSV-import voor auteurs.\nPDF-upload voor afrekeningen.'),
    ('Content', 'Evenementen, nieuws en vacatures\nbeheren vanuit het dashboard.'),
    ('Activiteiten', 'Live feed van logins en\nwijzigingsverzoeken.'),
]
for i, (title, desc) in enumerate(admin_feats):
    y = Inches(2.5 + i * 0.9)
    add_marker_dot(s, Inches(0.9), y + Inches(0.05), color=TEAL, size=0.08)
    add_text(s, Inches(1.15), y, Inches(3.3), Inches(0.25),
             title, sz=11, color=BLACK)
    add_text(s, Inches(1.15), y + Inches(0.28), Inches(3.3), Inches(0.6),
             desc, sz=9, color=LIGHT_GREY, spacing=2)


# ══════════════════════════════════════════════
# (Demo slide moved to after Pricing/Timeline)


# ══════════════════════════════════════════════
# SLIDE 10 — SCALABILITY
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(-0.5), 1.0)

add_subtitle_label(s, 'Schaalbaarheid')
add_title(s, 'E\u00E9n platform, drie organisaties', t=Inches(0.55))

orgs = [
    ('Noordhoff', '2.500 auteurs', 'Primair, voortgezet onderwijs\nen mbo/hbo in Nederland'),
    ('Liber', '~1.500 auteurs', 'Educatieve uitgeverij\nin Zweden'),
    ('Plantyn', '~1.000 auteurs', 'Educatieve uitgeverij\nin Belgi\u00EB'),
]
for i, (name, count, desc) in enumerate(orgs):
    l = Inches(0.9 + i * 4.0)

    # Marker (first one is highlighted)
    add_marker_dot(s, l, Inches(1.82), color=TEAL, size=0.14)
    add_text(s, l + Inches(0.3), Inches(1.75), Inches(3.2), Inches(0.35),
             name, sz=20, color=BLACK if i == 0 else BODY, font=FONT)
    add_text(s, l + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.3),
             count, sz=14, color=TEAL, font=FONT)
    add_text(s, l + Inches(0.3), Inches(2.6), Inches(3.2), Inches(1.0),
             desc, sz=12, color=LIGHT_GREY, spacing=5)

add_thin_line(s, Inches(0.9), Inches(4.1), Inches(11.5))

# White-label section
add_text(s, Inches(0.9), Inches(4.4), Inches(11), Inches(0.35),
         'White-label ready', sz=18, color=TEAL, font=FONT)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.3), Inches(1.5),
         'Elke organisatie krijgt een eigen huisstijl (logo, kleuren), eigen auteurs-database en eigen admin.\n'
         'De onderliggende technologie is identiek \u2014 updates worden \u00E9\u00E9n keer gebouwd en rollen uit naar alle drie.\n'
         'Geen extra ontwikkelkosten per organisatie, alleen configuratie en data-migratie.',
         sz=13, color=BODY, spacing=6)


# ══════════════════════════════════════════════
# SLIDE 11 — ROI
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.0), Inches(-0.5), 1.2)

add_subtitle_label(s, 'Businesscase')
add_title(s, 'Return on Investment', t=Inches(0.55))

# Stats row — clean numbers
roi_stats = [
    ('\u20AC10.000', 'Postkosten\nbespaard'),
    ('\u20AC5.000', 'Minder\nauteursvragen'),
    ('\u20AC2.500', 'Verwerking\nbespaard'),
    ('\u20AC17.500', 'Totale besparing\nper jaar'),
]
for i, (value, label) in enumerate(roi_stats):
    l = Inches(0.9 + i * 3.05)
    clr = TEAL
    add_text(s, l, Inches(1.6), Inches(2.7), Inches(0.6),
             value, sz=32, color=clr, font=FONT, align=PP_ALIGN.LEFT)
    add_text(s, l, Inches(2.2), Inches(2.7), Inches(0.5),
             label, sz=10, color=LIGHT_GREY, spacing=3)

# Highlight the total with an underline
add_thin_line(s, Inches(10.05), Inches(2.1), Inches(2.2), color=TEAL)

add_thin_line(s, Inches(0.9), Inches(3.0), Inches(11.5))

# ROI Table
tbl_shape = s.shapes.add_table(6, 3, Inches(0.9), Inches(3.3),
                                Inches(11.5), Inches(2.8))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(5.5)
tbl.columns[1].width = Inches(3)
tbl.columns[2].width = Inches(3)

rows = [
    ('', 'Huidig (per jaar)', 'Met portaal'),
    ('Fysieke post (print + porto + verzending)', '\u20AC10.000', '\u20AC0'),
    ('Afhandeling auteursvragen (125 uur)', '\u20AC5.000', '\u20AC500'),
    ('Fysieke verwerking (32 uur)', '\u20AC2.500', '\u20AC0'),
    ('Auteursbehoud & tevredenheid', 'Risico', 'Geborgd'),
    ('Totaal', '\u20AC17.500+', '\u20AC500'),
]
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = FONT
        p.font.bold = False

        if ri == 0:
            # Header row — teal background, white text
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL
            p.font.color.rgb = WHITE
        elif ri == len(rows) - 1:
            # Total row
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            p.font.color.rgb = TEAL
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            if ci == 1:
                p.font.color.rgb = ORANGE
            elif ci == 2:
                p.font.color.rgb = TEAL
            else:
                p.font.color.rgb = BODY

        if ci > 0:
            p.alignment = PP_ALIGN.CENTER

# Break-even visual — prominent
be_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.8))
be_bg.fill.solid(); be_bg.fill.fore_color.rgb = TEAL
be_bg.line.fill.background(); be_bg.adjustments[0] = 0.15

add_text(s, Inches(1.3), Inches(6.3), Inches(3.5), Inches(0.5),
         'Break-even: jaar 2', sz=18, color=WHITE, font=FONT)
add_text(s, Inches(5.0), Inches(6.35), Inches(7), Inches(0.5),
         'Vanaf jaar 2 bespaart het portaal meer dan het kost  (\u20AC17.500 besparing vs \u20AC15.000 licentie = +\u20AC2.500/jaar)',
         sz=11, color=RGBColor(200, 235, 225))


# ══════════════════════════════════════════════
# SLIDE 12 — PRICING
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(-0.5), Inches(-0.5), 1.0)

add_subtitle_label(s, 'Investering')
add_title(s, 'Wat kost het?', t=Inches(0.55))

# LEFT — Development cost
add_text(s, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.35),
         'Ontwikkeling & implementatie', sz=16, color=BLACK, font=FONT)
add_text(s, Inches(0.9), Inches(2.2), Inches(3), Inches(0.6),
         '\u20AC19.000', sz=40, color=TEAL, font=FONT)
add_text(s, Inches(3.6), Inches(2.5), Inches(2), Inches(0.25),
         'eenmalig', sz=12, color=LIGHT_GREY)
add_text(s, Inches(0.9), Inches(2.9), Inches(4.6), Inches(0.25),
         '200 uur \u00D7 \u20AC95/uur', sz=11, color=LIGHT_GREY)

dev_items = [
    'Volledig werkend portaal (auteur + admin)',
    'Database met Row Level Security',
    'Bulk import tooling (CSV + PDF)',
    'Tweetalig (NL/EN)',
    'Responsive (desktop, tablet, mobiel)',
    'Documentatie & overdracht',
]
for j, item in enumerate(dev_items):
    add_text(s, Inches(0.9), Inches(3.4 + j * 0.4), Inches(4.6), Inches(0.3),
             f'\u00B7   {item}', sz=11, color=BODY)

# Vertical separator
sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(6.4), Inches(1.7), Pt(1), Inches(4.5))
sep.fill.solid()
sep.fill.fore_color.rgb = BORDER_CLR
sep.line.fill.background()

# RIGHT — Annual license
add_text(s, Inches(6.9), Inches(1.7), Inches(5.2), Inches(0.35),
         'Jaarlijkse licentie', sz=16, color=BLACK, font=FONT)
add_text(s, Inches(6.9), Inches(2.2), Inches(3), Inches(0.6),
         '\u20AC5.000', sz=40, color=TEAL, font=FONT)
add_text(s, Inches(9.2), Inches(2.5), Inches(2.5), Inches(0.25),
         '/jaar/organisatie', sz=12, color=LIGHT_GREY)
add_text(s, Inches(6.9), Inches(2.9), Inches(4.6), Inches(0.25),
         '3 organisaties = \u20AC15.000/jaar', sz=11, color=LIGHT_GREY)

license_items = [
    'Hosting & infrastructuur',
    'Onderhoud, updates & bugfixes',
    'Support (mail + telefoon)',
    'Branding per organisatie (white-label)',
    'Data-migratie ondersteuning',
    'Nieuwe features op aanvraag',
]
for j, item in enumerate(license_items):
    add_text(s, Inches(6.9), Inches(3.4 + j * 0.4), Inches(4.6), Inches(0.3),
             f'\u00B7   {item}', sz=11, color=BODY)

# Bottom summary
add_thin_line(s, Inches(0.9), Inches(6.2), Inches(11.5), color=TEAL)
add_text(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.35),
         'Jaar 1: \u20AC34.000  \u00B7  Besparing: \u20AC17.500  \u00B7  Netto: \u20AC16.500  |  '
         'Vanaf jaar 2: \u20AC15.000 vs \u20AC17.500 besparing = netto positief',
         sz=12, color=TEAL, align=PP_ALIGN.CENTER, font=FONT)


# ══════════════════════════════════════════════
# SLIDE 13 — TIMELINE
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(5.8), 1.8)

add_subtitle_label(s, 'Implementatie')
add_title(s, 'Tijdlijn naar productie', t=Inches(0.55))

phases = [
    ('Week 1\u20132', 'Setup', 'Contract, toegang,\ndata-inventarisatie'),
    ('Week 3\u20134', 'Data-migratie', 'Auteursdata importeren,\nPDFs uploaden'),
    ('Week 5\u20136', 'Testen', 'Pilot met 10 auteurs,\nfeedback verwerken'),
    ('Week 7', 'Go-live', 'Alle auteurs ontvangen\nlogin per email'),
]
for i, (week, phase, desc) in enumerate(phases):
    l = Inches(0.9 + i * 3.05)

    # Phase number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, l, Inches(1.8),
                                 Inches(0.4), Inches(0.4))
    fill_clr = TEAL if i == 3 else WHITE
    txt_clr = WHITE if i == 3 else TEAL
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_clr
    if i < 3:
        circle.line.color.rgb = TEAL
        circle.line.width = Pt(1.5)
    else:
        circle.line.fill.background()
    add_text(s, l + Inches(0.02), Inches(1.82), Inches(0.36), Inches(0.36),
             str(i + 1), sz=14, color=txt_clr, align=PP_ALIGN.CENTER, font=FONT)

    # Connecting line between circles
    if i < 3:
        conn = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   l + Inches(0.45), Inches(1.98),
                                   Inches(2.55), Pt(2))
        conn.fill.solid()
        conn.fill.fore_color.rgb = BORDER_CLR
        conn.line.fill.background()

    # Week label
    add_text(s, l, Inches(2.35), Inches(2.5), Inches(0.25),
             week, sz=10, color=TEAL, font=FONT)
    # Phase name
    add_text(s, l, Inches(2.65), Inches(2.5), Inches(0.35),
             phase, sz=18, color=BLACK, font=FONT)
    # Description
    add_text(s, l, Inches(3.1), Inches(2.5), Inches(1.0),
             desc, sz=12, color=BODY, spacing=5)

add_thin_line(s, Inches(0.9), Inches(4.8), Inches(11.5))

# Bottom text
add_text(s, Inches(0.9), Inches(5.1), Inches(11.3), Inches(1.0),
         'Het portaal is al gebouwd en getest. De enige stap naar productie is het importeren van echte data.\n'
         'Bij akkoord kan het portaal binnen 7 weken live zijn voor alle 2.500 auteurs.',
         sz=14, color=BODY, spacing=6)

# Handwritten note
add_text(s, Inches(7.5), Inches(6.0), Inches(5), Inches(0.5),
         'klaar om te lanceren!', sz=22, color=BODY, font=FONT_HAND)


# ══════════════════════════════════════════════
# SLIDE 14 — DEMO (Green background, after business slides)
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK_GREEN)

for (pl, pt, pw, ph) in [
    (Inches(0.5), Inches(0.4), Inches(2.0), Inches(0.3)),
    (Inches(10.5), Inches(6.6), Inches(2.5), Inches(0.4)),
    (Inches(11.8), Inches(0.3), Inches(1.2), Inches(0.2)),
]:
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pl, pt, pw, ph)
    pill.fill.solid(); pill.fill.fore_color.rgb = WHITE_SHAPE
    pill.line.fill.background(); pill.adjustments[0] = 0.5

add_text(s, Inches(0), Inches(2.8), W, Inches(1.2),
         'Demo', sz=72, color=WHITE, align=PP_ALIGN.CENTER, font=FONT)


# ══════════════════════════════════════════════
# SLIDE 15 — CTA with concrete next step
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK_GREEN)

add_text(s, Inches(0), Inches(1.4), W, Inches(1.0),
         'Volgende stap', sz=48, color=WHITE,
         align=PP_ALIGN.CENTER, font=FONT)

add_text(s, Inches(2.5), Inches(2.6), Inches(8.3), Inches(0.6),
         'Het portaal is gebouwd en klaar voor productie.',
         sz=17, color=RGBColor(200, 235, 225), align=PP_ALIGN.CENTER)

# Concrete next steps
steps = [
    ('1', 'Akkoord op ontwikkelbudget'),
    ('2', 'Toegang tot auteursdata en systemen'),
    ('3', 'Data-migratie en configuratie (5 weken)'),
    ('4', 'Pilot met 10 auteurs (1 week)'),
    ('5', 'Go-live: alle 2.500 auteurs ontvangen een login'),
]
for i, (num, step) in enumerate(steps):
    y = Inches(3.5 + i * 0.5)
    # Number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), y, Inches(0.3), Inches(0.3))
    circ.fill.solid(); circ.fill.fore_color.rgb = WHITE_SHAPE; circ.line.fill.background()
    add_text(s, Inches(4.2), y + Inches(0.02), Inches(0.3), Inches(0.26),
             num, sz=11, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.7), y + Inches(0.02), Inches(5), Inches(0.3),
             step, sz=14, color=WHITE)

add_text(s, Inches(0), Inches(6.3), W, Inches(0.3),
         'Patrick Jeeninga  \u00B7  patrick@noordhoff.nl',
         sz=12, color=RGBColor(140, 195, 180), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# SLIDE 16 — APPENDIX: SECURITY
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(12.2), Inches(-0.5), 1.0)

add_text(s, Inches(0.9), Inches(0.3), Inches(3), Inches(0.25),
         'APPENDIX', sz=9, color=LIGHT_GREY, font=FONT)
add_title(s, 'Databeveiliging', t=Inches(0.5))

# Security layers — left side
layers = [
    ('Authenticatie', 'Inloggen via versleutelde wachtwoorden (bcrypt).\nSessie-tokens verlopen automatisch.'),
    ('Row Level Security', 'Database dwingt af dat elke auteur alleen eigen\ndata ziet. Zelfs bij directe API-toegang.'),
    ('Versleutelde verbinding', 'Alle communicatie via HTTPS/TLS.\nData in transit is altijd versleuteld.'),
    ('Opslag-beveiliging', 'PDFs in priv\u00E9 bucket. Download alleen via\ntijdelijke signed URLs (verlopen na 5 min).'),
    ('BSN-maskering', 'BSN wordt gemaskeerd weergegeven (\u2022\u2022\u2022\u2022\u20226789).\nPas zichtbaar na expliciete actie.'),
]
for i, (lbl, desc) in enumerate(layers):
    t = Inches(1.6 + i * 0.95)
    add_marker_dot(s, Inches(1.0), t + Inches(0.06), color=TEAL)
    add_text(s, Inches(1.3), t, Inches(2.2), Inches(0.3),
             lbl, sz=13, color=BLACK, font=FONT)
    add_text(s, Inches(1.3), t + Inches(0.35), Inches(5.5), Inches(0.6),
             desc, sz=11, color=BODY, spacing=3)

# Right side — comparison
# Risks of physical mail (with orange markers)
add_text(s, Inches(8.0), Inches(1.6), Inches(4), Inches(0.3),
         "Risico's fysieke post", sz=13, color=ORANGE, font=FONT)

risks = [
    'Brieven met financi\u00EBle data raken zoek',
    'Geen controle wie de brief opent',
    'Geen audit trail',
    "Kopie\u00EBn op bureaus",
]
for j, r in enumerate(risks):
    add_marker_dot(s, Inches(8.0), Inches(2.15 + j * 0.4), color=ORANGE, size=0.08)
    add_text(s, Inches(8.2), Inches(2.08 + j * 0.4), Inches(4), Inches(0.3),
             r, sz=11, color=BODY)

add_thin_line(s, Inches(8.0), Inches(3.8), Inches(4.3))

# Guarantees of portal (with teal markers)
add_text(s, Inches(8.0), Inches(4.1), Inches(4), Inches(0.3),
         'Waarborgen portaal', sz=13, color=TEAL, font=FONT)

guarantees = [
    'Toegang alleen na inloggen',
    'Auteur ziet uitsluitend eigen data (RLS)',
    'Elke login wordt gelogd (audit trail)',
    'Admin kan activiteit monitoren',
    'Geen fysieke documenten',
]
for j, r in enumerate(guarantees):
    add_marker_dot(s, Inches(8.0), Inches(4.65 + j * 0.38), color=TEAL, size=0.08)
    add_text(s, Inches(8.2), Inches(4.58 + j * 0.38), Inches(4), Inches(0.3),
             r, sz=11, color=BODY)


# ══════════════════════════════════════════════
# SLIDE 17 — APPENDIX: AVG COMPLIANCE
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_teal_square(s, Inches(-0.5), Inches(5.8), 1.2)

add_text(s, Inches(0.9), Inches(0.3), Inches(3), Inches(0.25),
         'APPENDIX', sz=9, color=LIGHT_GREY, font=FONT)
add_title(s, 'AVG-compliance', t=Inches(0.5))

avg = [
    ('Recht op inzage (Art. 15)',
     'Auteurs kunnen 24/7 eigen gegevens,\ncontracten en afrekeningen inzien.',
     'Geen formeel verzoek nodig \u2014\nhet portaal biedt dit standaard.'),
    ('Recht op rectificatie (Art. 16)',
     'Wijzigingsverzoek indienen via portaal\n(adres, telefoon, IBAN).',
     'Admin keurt goed of wijst af.\nVolledig traceerbaar.'),
    ('Dataminimalisatie (Art. 5)',
     'Alleen relevante gegevens getoond.\nBSN is standaard gemaskeerd.',
     'Geen onnodige data-opslag\nof verspreiding.'),
    ('Beveiliging (Art. 32)',
     'HTTPS, gehashte wachtwoorden, RLS,\nsigned URLs voor downloads.',
     'Technische maatregelen zijn\ningebouwd in de architectuur.'),
    ('Verantwoording (Art. 5)',
     'Login-logging en audit trail op alle\nwijzigingsverzoeken met timestamps.',
     'Bij een audit exact traceerbaar\nwie wat heeft gedaan.'),
]

# Column headers
add_text(s, Inches(0.9), Inches(1.3), Inches(3), Inches(0.3),
         'Principe', sz=10, color=TEAL, font=FONT)
add_text(s, Inches(4.2), Inches(1.3), Inches(3.5), Inches(0.3),
         'Hoe het portaal dit invult', sz=10, color=TEAL, font=FONT)
add_text(s, Inches(8.5), Inches(1.3), Inches(3.5), Inches(0.3),
         'Toelichting', sz=10, color=TEAL, font=FONT)

add_thin_line(s, Inches(0.9), Inches(1.6), Inches(11.5))

for i, (principle, how, note) in enumerate(avg):
    t = Inches(1.8 + i * 1.0)

    add_marker_dot(s, Inches(0.9), t + Inches(0.06), color=TEAL, size=0.1)
    add_text(s, Inches(1.15), t, Inches(2.8), Inches(0.65),
             principle, sz=11, color=BLACK, font=FONT)
    add_text(s, Inches(4.2), t, Inches(3.8), Inches(0.65),
             how, sz=10, color=BODY, spacing=3)
    add_text(s, Inches(8.5), t, Inches(3.8), Inches(0.65),
             note, sz=10, color=LIGHT_GREY, spacing=3)

    # Separator between rows
    if i < len(avg) - 1:
        add_thin_line(s, Inches(0.9), t + Inches(0.82), Inches(11.5))

# Bottom summary
add_text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.3),
         'Het portaal vervangt fysieke brieven \u2014 dit vermindert het risico op '
         'datalekken en verbetert de compliance-positie.',
         sz=12, color=TEAL, align=PP_ALIGN.CENTER, font=FONT)


# ══════════════════════════════════════════════
# APPENDIX 3 — WORK LOG (moved from main deck)
# ══════════════════════════════════════════════

s = prs.slides.add_slide(LY_BLANK)

add_text(s, Inches(0.9), Inches(0.3), Inches(3), Inches(0.25),
         'APPENDIX', sz=9, color=LIGHT_GREY, font=FONT)
add_title(s, 'Urenverantwoording \u2014 200 uur', t=Inches(0.5))

log = [
    ('UX research & wireframes', 16),
    ('Huisstijl & design system', 12),
    ('Login & authenticatie', 10),
    ('Dashboard framework & routing', 14),
    ('7 auteur-tabbladen bouwen', 28),
    ('PDF preview & generatie', 10),
    ('Admin dashboard & CRUD', 22),
    ('Supabase backend & RLS', 18),
    ('Bulk import tools (CSV + PDF)', 12),
    ('Responsive design (mobile/tablet)', 14),
    ('i18n, dark mode, guided tour', 12),
    ("Publieke website (5 pagina's)", 16),
    ('Testing, QA & optimalisatie', 16),
]
max_hours = max(h for _, h in log)
for i, (task, hours) in enumerate(log):
    col, row = i % 2, i // 2
    l = Inches(0.9 + col * 6.1)
    t = Inches(1.7 + row * 0.6)
    add_text(s, l, t, Inches(4.0), Inches(0.35), task, sz=11, color=BODY)
    add_text(s, l + Inches(4.0), t, Inches(1.5), Inches(0.35),
             f'{hours} uur', sz=11, color=TEAL, align=PP_ALIGN.RIGHT, font=FONT)
    bar_width = Inches(5.2 * hours / max_hours)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              l, t + Inches(0.35), bar_width, Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background(); bar.adjustments[0] = 0.5

add_thin_line(s, Inches(0.9), Inches(6.1), Inches(11.5), color=TEAL)
add_text(s, Inches(0.9), Inches(6.2), Inches(5), Inches(0.35),
         'TOTAAL', sz=14, color=TEAL, font=FONT)
add_text(s, Inches(10.2), Inches(6.2), Inches(2.1), Inches(0.35),
         '200 uur', sz=14, color=TEAL, align=PP_ALIGN.RIGHT, font=FONT)


# ══════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════

out = '/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/Auteursportaal_Pitch.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
