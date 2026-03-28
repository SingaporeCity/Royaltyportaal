#!/usr/bin/env python3
"""Full deep inspection of every slide in the Noordhoff template"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn

prs = Presentation('/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/Template richtlijnen - 2026.pptx')

print(f"Slides: {len(prs.slides)}")
print(f"Size: {prs.slide_width/914400:.2f}\" x {prs.slide_height/914400:.2f}\"")

for si, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {si+1}: layout='{slide.slide_layout.name}'")
    print(f"{'='*70}")

    for shape in slide.shapes:
        name = shape.name
        st = shape.shape_type
        l = shape.left/914400
        t = shape.top/914400
        w = shape.width/914400
        h = shape.height/914400

        # Get fill color
        fill_str = ""
        el = shape._element
        for srgb in el.iter(qn('a:srgbClr')):
            fill_str += f" #{srgb.get('val')}"

        # Get text
        txt = ""
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text[:50].replace('\n', '|')

        # Shape type name
        type_name = str(st).split('(')[0].split('.')[-1] if '.' in str(st) else str(st)

        print(f"  [{name}] type={type_name}")
        print(f"    pos=({l:.1f}\",{t:.1f}\") size=({w:.1f}\"x{h:.1f}\"){fill_str}")
        if txt:
            print(f"    text: \"{txt}\"")

        # Inspect group children
        if shape.shape_type == 6:  # GROUP
            try:
                for child in shape.shapes:
                    cl = child.left/914400
                    ct = child.top/914400
                    cw = child.width/914400
                    ch = child.height/914400
                    child_colors = ""
                    for srgb in child._element.iter(qn('a:srgbClr')):
                        child_colors += f" #{srgb.get('val')}"
                    child_txt = ""
                    if hasattr(child, 'text') and child.text:
                        child_txt = child.text[:30]
                    print(f"      └─ [{child.name}] ({cw:.2f}\"x{ch:.2f}\"){child_colors} {child_txt}")
            except:
                pass

# Also check what the green layouts look like
print(f"\n{'='*70}")
print("LAYOUT BACKGROUNDS")
print(f"{'='*70}")
for i, layout in enumerate(prs.slide_layouts):
    bg = layout.background
    colors = ""
    for srgb in layout.element.iter(qn('a:srgbClr')):
        colors += f" #{srgb.get('val')}"
    print(f"  Layout {i} '{layout.name}': {colors[:60] if colors else 'theme colors'}")
