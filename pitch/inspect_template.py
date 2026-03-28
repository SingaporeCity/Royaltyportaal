#!/usr/bin/env python3
from pptx import Presentation

prs = Presentation('/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/Template richtlijnen - 2026.pptx')

for si, slide in enumerate(prs.slides):
    print(f"\n--- Slide {si+1}: '{slide.slide_layout.name}' ---")
    for shape in slide.shapes:
        txt = shape.text[:80].replace('\n', ' | ') if hasattr(shape, 'text') and shape.text else ''
        print(f"  {shape.name}: {txt}")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs[:2]:
                for r in p.runs[:1]:
                    try: c = str(r.font.color.rgb)
                    except: c = 'theme'
                    print(f"    font='{r.font.name}' sz={r.font.size} bold={r.font.bold} color={c}")
