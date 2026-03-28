#!/usr/bin/env python3
"""Deep inspect Noordhoff template — extract exact colors, shapes, markers, decorative elements"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
import json

prs = Presentation('/Users/patrickjeeninga/Coding/Royaltyportaal/pitch/Template richtlijnen - 2026.pptx')

# Check theme colors
print("=== THEME COLORS ===")
theme = prs.slide_masters[0].slide_layouts[0].slide_master.element
for clr in theme.iter(qn('a:srgbClr')):
    print(f"  srgbClr: #{clr.get('val')}")
for clr in theme.iter(qn('a:schemeClr')):
    print(f"  schemeClr: {clr.get('val')}")

# Inspect slide 1 deeply for markers/decorative shapes
print("\n=== SLIDE 1 DEEP (markers, groups) ===")
slide = prs.slides[0]
for shape in slide.shapes:
    print(f"\n  Shape: '{shape.name}' type={shape.shape_type}")
    print(f"    pos: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\") size: ({shape.width/914400:.2f}\"x{shape.height/914400:.2f}\")")

    # Check fill colors
    if hasattr(shape, 'fill'):
        fill = shape.fill
        try:
            if fill.type is not None:
                print(f"    fill type: {fill.type}")
        except: pass

    # For groups, inspect children
    if shape.shape_type == 6:  # GROUP
        try:
            for child in shape.shapes:
                print(f"      Child: '{child.name}' type={child.shape_type}")
                print(f"        pos: ({child.left/914400:.2f}\", {child.top/914400:.2f}\") size: ({child.width/914400:.2f}\"x{child.height/914400:.2f}\")")
                # Check for fill color
                el = child._element
                for srgb in el.iter(qn('a:srgbClr')):
                    print(f"        color: #{srgb.get('val')}")
        except Exception as e:
            print(f"      Error: {e}")

# Inspect slide 4 (blank with rectangles - content layout)
print("\n=== SLIDE 4 (content layout) ===")
slide = prs.slides[3]
for shape in slide.shapes:
    print(f"\n  Shape: '{shape.name}' type={shape.shape_type}")
    print(f"    pos: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\") size: ({shape.width/914400:.2f}\"x{shape.height/914400:.2f}\")")
    el = shape._element
    for srgb in el.iter(qn('a:srgbClr')):
        print(f"    color: #{srgb.get('val')}")
    if hasattr(shape, 'text') and shape.text:
        print(f"    text: {shape.text[:60]}")

# Inspect slide 10 (Noordhoff in cijfers - with stat groups)
print("\n=== SLIDE 10 (stats) ===")
slide = prs.slides[9]
for shape in slide.shapes:
    print(f"\n  Shape: '{shape.name}' type={shape.shape_type}")
    if hasattr(shape, 'text') and shape.text:
        print(f"    text: {shape.text[:60]}")
    if shape.shape_type == 6:
        try:
            for child in shape.shapes:
                print(f"      Child: '{child.name}' t={child.shape_type}")
                if hasattr(child, 'text') and child.text:
                    print(f"        text: {child.text[:60]}")
                el = child._element
                for srgb in el.iter(qn('a:srgbClr')):
                    print(f"        color: #{srgb.get('val')}")
        except Exception as e:
            print(f"      Error: {e}")

# Inspect slide 9 (section header with shapes)
print("\n=== SLIDE 9 (section with markers) ===")
slide = prs.slides[8]
for shape in slide.shapes:
    print(f"\n  Shape: '{shape.name}' type={shape.shape_type}")
    print(f"    pos: ({shape.left/914400:.2f}\", {shape.top/914400:.2f}\") size: ({shape.width/914400:.2f}\"x{shape.height/914400:.2f}\")")
    el = shape._element
    for srgb in el.iter(qn('a:srgbClr')):
        print(f"    color: #{srgb.get('val')}")
    if hasattr(shape, 'text') and shape.text:
        print(f"    text: {shape.text[:60]}")
